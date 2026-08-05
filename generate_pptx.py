import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(15, 23, 42)
DARK_NAVY = RGBColor(9, 13, 22)
GOLD = RGBColor(212, 175, 55)
GREEN = RGBColor(0, 135, 81)
WHITE = RGBColor(255, 255, 255)
TEXT_MUTED = RGBColor(148, 163, 184)
CARD_BG = RGBColor(30, 41, 59)

def add_header(slide, title_text, category_text="DIGITs NIGERIA ELECTION WATCH — COMMERCIAL PROJECT INVOICE"):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_NAVY
    bg.line.fill.background()
    
    cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf = cat.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = category_text
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    t = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

# Slide 1: Cover / Header
blank_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_layout)

bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = DARK_NAVY
bg1.line.fill.background()

tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.8))
tf1 = tb1.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = "DIGITs NIGERIA ELECTION WATCH"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = GOLD

p2 = tf1.add_paragraph()
p2.text = "Expanded Commercial Project Invoice"
p2.font.size = Pt(38)
p2.font.bold = True
p2.font.color.rgb = WHITE
p2.space_after = Pt(8)

p3 = tf1.add_paragraph()
p3.text = "Full Scalable Web, Mobile (Android & iOS), and 10M Concurrent User Backend Architecture"
p3.font.size = Pt(15)
p3.font.color.rgb = TEXT_MUTED
p3.space_after = Pt(20)

p4 = tf1.add_paragraph()
p4.text = "REVISED TOTAL VALUE: ₦58,000,000.00 (FIFTY-EIGHT MILLION NAIRA)"
p4.font.size = Pt(22)
p4.font.bold = True
p4.font.color.rgb = GREEN

# Details Box
box1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.6))
box1.fill.solid()
box1.fill.fore_color.rgb = CARD_BG
box1.line.color.rgb = GOLD
tf = box1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Invoice #: INV-2026-DIGEO-058M  |  Date: August 5, 2026  |  Issuer: WYN-Tech Systems Ltd. (SirHope of WYN-Tech)"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.space_after = Pt(6)
p2 = tf.add_paragraph()
p2.text = "Target Scope: 176,846 Polling Units  |  36 States & FCT  |  10,000,000 Peak Active Citizen Viewers Load Rating"
p2.font.size = Pt(11)
p2.font.color.rgb = TEXT_MUTED

# Slide 2: Expanded Technical Architecture Scope
slide2 = prs.slides.add_slide(blank_layout)
add_header(slide2, "Expanded Scope & High-Concurrency Architecture Specs")

card_data = [
    ("Web & Control Center", "• React 19, Vite, TanStack Router\n• OKLCH Brand Design System\n• 1-to-6 Split Screen WebRTC Grid\n• 11-Screen Command Center\n• Immutable Audit Trail"),
    ("Cross-Platform Mobile Apps", "• Native iOS & Android Builds\n• Forced GPS Location Gate\n• 2-Min Live Camera Recorder\n• Gallery Upload Blocking\n• Cryptographic NIN Identity Hash"),
    ("10M User Backend & Load", "• Supabase Auto-Scaling Cluster\n• LiveKit WebRTC Intake Server\n• Redis Caching & DB Sharding\n• k6 / Locust 10M Load Tested\n• Multi-Region DDoS Shield")
]

for i, (ctitle, cdesc) in enumerate(card_data):
    left = Inches(0.8 + i * 4.0)
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(5.0))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = GOLD if i==2 else NAVY
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ctitle
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = cdesc
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = WHITE

# Slide 3: Revised Commercial Invoice Line Items (₦58,000,000)
slide3 = prs.slides.add_slide(blank_layout)
add_header(slide3, "Itemized Financial Breakdown & Deliverables (₦58M Total)")

items = [
    ("1. Web Application & Control Center", "React 19, Vite, TanStack Router, OKLCH design system, 6-split WebRTC Live Video Grid, i-Witness evidence queue, 11-screen Command Center dashboard.", "₦ 12,500,000"),
    ("2. Cross-Platform Mobile Apps (Android & iOS)", "Native iOS & Android apps. On-the-spot 2-min camera capture, forced GPS check, NIN hash verification, gallery upload blocking, offline sync.", "₦ 15,000,000"),
    ("3. 10M Concurrent User High-Availability Backend", "Supabase Enterprise Cluster, LiveKit WebRTC Intake Engine, PostgreSQL auto-sharding, Read Replicas, Redis Caching, Edge Functions for 10M active users.", "₦ 17,500,000"),
    ("4. DIGEO Accreditation Academy & Verification", "6 self-paced interactive modules, 24 EC8A & Electoral Law assessment items, automated QR certificate generation, RLS role-based access control.", "₦ 5,500,000"),
    ("5. Stress Testing, QA, Security Audit & SLA", "k6/Locust load testing simulating 10,000,000 active viewers, failover redundancy, multi-region DDoS mitigation, Penetration Audit & 12M SLA.", "₦ 7,500,000")
]

rows = 6
cols = 3
table_shape = slide3.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.7))
table = table_shape.table
table.columns[0].width = Inches(3.6)
table.columns[1].width = Inches(5.6)
table.columns[2].width = Inches(2.5)

headers = ["Module / Deliverable", "Expanded Technical Specifications", "Amount (NGN)"]
for col_idx, htext in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.text = htext
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GOLD

for row_idx, (m_title, m_desc, m_amt) in enumerate(items, start=1):
    cell_0 = table.cell(row_idx, 0)
    cell_1 = table.cell(row_idx, 1)
    cell_2 = table.cell(row_idx, 2)
    
    for c in [cell_0, cell_1, cell_2]:
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        
    p0 = cell_0.text_frame.paragraphs[0]
    p0.text = m_title
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = WHITE
    
    p1 = cell_1.text_frame.paragraphs[0]
    p1.text = m_desc
    p1.font.size = Pt(9)
    p1.font.color.rgb = TEXT_MUTED
    
    p2 = cell_2.text_frame.paragraphs[0]
    p2.text = m_amt
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

# Total Banner
total_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6))
total_box.fill.solid()
total_box.fill.fore_color.rgb = CARD_BG
total_box.line.color.rgb = GOLD
tf = total_box.text_frame
p = tf.paragraphs[0]
p.text = "TOTAL REVISED VALUE: ₦ 58,000,000.00 (FIFTY-EIGHT MILLION NAIRA)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = GREEN
p.alignment = PP_ALIGN.CENTER

# Slide 4: Terms & Sign-off
slide4 = prs.slides.add_slide(blank_layout)
add_header(slide4, "Payment Schedule, Milestone Breakdown & Support SLA")

terms_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
terms_card.fill.solid()
terms_card.fill.fore_color.rgb = CARD_BG
terms_card.line.color.rgb = GOLD
tf = terms_card.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "REVISED PAYMENT MILESTONES (TOTAL: ₦58M)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = GOLD
p.space_after = Pt(6)

milestones = [
    "• 40% Mobilization Deposit (₦ 23,200,000.00) — Architecture Setup, Web & Mobile Core Build, 10M Cluster Setup",
    "• 40% Beta Delivery & Load Test Signoff (₦ 23,200,000.00) — 10M Load Test Pass, Store Submission & Audit Signoff",
    "• 20% Production Handover (₦ 11,600,000.00) — Election Day Live Control Center Standby & Final Acceptance"
]
for m in milestones:
    p = tf.add_paragraph()
    p.text = m
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.space_after = Pt(5)

p = tf.add_paragraph()
p.text = "\nBANKING & ENGINEERING SLA DETAILS"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = GOLD
p.space_after = Pt(6)

details = [
    "• Corporate Name: WYN-Tech Systems Ltd.",
    "• Lead Architect: SirHope of WYN-Tech (wyntech.ng@gmail.com)",
    "• Bank: Zenith Bank Plc / Access Bank Nigeria | Account #: 1012345678",
    "• Support SLA: 12 months full technical support, 15-minute emergency response, 99.99% uptime guarantee."
]
for d in details:
    p = tf.add_paragraph()
    p.text = d
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.space_after = Pt(4)

pptx_filename = "c:/Users/danho/Desktop/digits-naija-election-watch/DIGITs_Election_Watch_Invoice.pptx"
prs.save(pptx_filename)
print("PPTX generated successfully:", pptx_filename)
